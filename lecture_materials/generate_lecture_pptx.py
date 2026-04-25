"""Generate a 17-slide lecture presentation for a Texas Bridge Risk Assessment workshop."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path(__file__).resolve().parent
PPTX_PATH = OUTPUT_DIR / "lecture_bridge_risk.pptx"

FIG_DIR = Path(__file__).resolve().parent.parent / "bridge_risk_demo" / "outputs"
FIG1 = FIG_DIR / "fig1_risk_map.png"
FIG2 = FIG_DIR / "fig2_age_condition.png"
FIG4 = FIG_DIR / "fig4_county_risk.png"

DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT = RGBColor(0xD4, 0x6B, 0x08)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MED_GREY = RGBColor(0x66, 0x66, 0x66)
TABLE_HDR = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ALT = RGBColor(0xE8, 0xEE, 0xF4)

SW = Inches(13.333)
SH = Inches(7.5)


def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def _tb(slide, l, t, w, h, text, sz=18, bold=False, color=DARK_GREY,
        align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return box


def _bullets(slide, l, t, w, h, items, sz=18, color=DARK_GREY, sp=Pt(10), bold=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = sp
    return box


def _header(slide, title):
    _rect(slide, Inches(0), Inches(0), SW, Inches(1.15), DARK_BLUE)
    _tb(slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
        title, sz=30, bold=True, color=WHITE)


def _snum(slide, n):
    _tb(slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35),
        str(n), sz=10, color=MED_GREY, align=PP_ALIGN.RIGHT)


def _notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def _table(slide, data, left, top, width, height, col_widths=None):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = ts.table
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = w
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.name = "Calibri"
                if r == 0:
                    p.font.bold = True
                    p.font.color.rgb = WHITE
                else:
                    p.font.color.rgb = DARK_GREY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HDR
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT
    return ts


def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    bl = prs.slide_layouts[6]

    # ── Slide 1: Title ──
    s = prs.slides.add_slide(bl)
    _bg(s, DARK_BLUE)
    _tb(s, Inches(1), Inches(1.3), Inches(11.3), Inches(1.5),
        "Bridge Risk Assessment", sz=42, bold=True, color=WHITE)
    _tb(s, Inches(1), Inches(2.8), Inches(11.3), Inches(1.0),
        "A Data-Driven Workshop Using the Texas National Bridge Inventory",
        sz=24, color=RGBColor(0xB0, 0xCC, 0xE6))
    _tb(s, Inches(1), Inches(4.6), Inches(11.3), Inches(0.5),
        "Civil & Natural Resources Engineering", sz=18, color=RGBColor(0xB0, 0xCC, 0xE6))
    _tb(s, Inches(1), Inches(5.1), Inches(11.3), Inches(0.5),
        "University of Canterbury", sz=16, color=RGBColor(0x90, 0xAA, 0xCC))
    _rect(s, Inches(0), Inches(6.7), SW, Inches(0.8), RGBColor(0x14, 0x2D, 0x4A))
    _tb(s, Inches(1), Inches(6.78), Inches(11.3), Inches(0.5),
        "Data Analysis Workshop  |  2025  |  Interactive Session",
        sz=14, color=RGBColor(0x90, 0xAA, 0xCC))
    _notes(s,
        "Welcome to the Bridge Risk Assessment workshop. Today we will work through a "
        "complete data analysis pipeline: loading real infrastructure data from the FHWA "
        "National Bridge Inventory, computing composite risk scores, and exploring the "
        "results interactively. By the end of this session, you will have hands-on "
        "experience with an open dataset that is used by every state transport agency in "
        "the United States. The skills transfer directly to asset management practice.")

    # ── Slide 2: Real-World Failures ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Why Bridge Risk Matters")
    _bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), [
        "I-35W Mississippi River bridge collapse (2007) -- 13 fatalities, rated 'structurally deficient'",
        "Fern Hollow Bridge, Pittsburgh (2022) -- collapsed hours before presidential visit",
        "Morandi Bridge, Genoa (2018) -- 43 deaths; cable-stayed span failed during routine traffic",
        "Common thread: deterioration was documented, but not acted on in time",
    ], sz=21, sp=Pt(20))
    _tb(s, Inches(0.8), Inches(5.9), Inches(11.5), Inches(0.5),
        "Question: How do we move from documenting deficiency to prioritising intervention?",
        sz=18, bold=True, color=MID_BLUE)
    _snum(s, 2)
    _notes(s,
        "These three bridge failures illustrate why risk assessment matters. In each case, "
        "inspection data showed deterioration before collapse. The I-35W bridge had been "
        "rated structurally deficient for over a decade. The lesson is not that we lack data "
        "-- the NBI collects detailed condition ratings on every public bridge. The problem "
        "is translating that data into prioritised action. With tens of thousands of bridges "
        "in a single state, agencies need a systematic framework to decide where to allocate "
        "limited inspection and rehabilitation budgets. That is exactly what we build today.")

    # ── Slide 3: The Prioritisation Problem ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "The Prioritisation Problem")
    _bullets(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(4.0), [
        "56,951 structures in Texas alone",
        "Limited annual maintenance budgets",
        "Condition ratings are necessary but insufficient",
        "Age and traffic exposure amplify structural risk",
        "Need: a composite metric for ranking",
    ], sz=20, sp=Pt(16))
    _table(s, [
        ["Metric", "What It Captures", "Limitation"],
        ["Condition rating", "Current physical state", "Snapshot; no trend"],
        ["Bridge age", "Cumulative deterioration", "Ignores maintenance"],
        ["Traffic (ADT)", "Loading & consequences", "Not structural"],
        ["Composite score", "All three factors", "Weight selection"],
    ], Inches(6.8), Inches(1.6), Inches(6.2), Inches(3.0),
       col_widths=[Inches(1.8), Inches(2.4), Inches(2.0)])
    _snum(s, 3)
    _notes(s,
        "Texas has the largest bridge inventory in the United States. Inspecting and "
        "maintaining nearly 57,000 structures requires systematic prioritisation. Condition "
        "ratings alone tell you what is wrong now, but not how urgently it needs attention. "
        "A bridge rated 5 out of 9 on a quiet rural road is very different from one rated 5 "
        "on an interstate carrying 200,000 vehicles per day. Age matters because older bridges "
        "have accumulated more fatigue cycles and environmental exposure. The composite score "
        "we develop today combines all three dimensions into a single ranking metric.")

    # ── Slide 4: Workshop Roadmap ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Workshop Roadmap")
    steps = [
        "1.  Understand the NBI dataset and key variables",
        "2.  Learn the three risk components: condition, age, traffic",
        "3.  Compute a weighted composite risk score",
        "4.  Classify bridges into risk tiers",
        "5.  Interpret results: maps, distributions, and county patterns",
        "6.  Explore interactively: adjust weights, compare scenarios",
    ]
    _bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), steps, sz=22, sp=Pt(18))
    _snum(s, 4)
    _notes(s,
        "Here is our roadmap. We begin with the data -- understanding what the NBI records "
        "and how to clean it. Then we define three risk sub-scores, each normalised between "
        "zero and one. We combine them using a weighted formula and classify every bridge "
        "into Low, Moderate, High, or Critical tiers. The bulk of the session covers "
        "interpreting the results through maps, scatter plots, and distributions. Finally, "
        "you will use the interactive explorer tool to adjust the weights yourself and see "
        "how the risk landscape changes. Each step builds on the last.")

    # ── Slide 5: NBI Dataset Overview ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "The National Bridge Inventory")
    _bullets(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), [
        "Maintained by FHWA; updated annually",
        "Every public highway bridge in the U.S.",
        "123 columns per record",
        "Texas 2025: 56,951 structures",
        "Publicly available (open data)",
    ], sz=19, sp=Pt(12))
    _table(s, [
        ["NBI Column", "Description", "Type"],
        ["STRUCTURE_NUMBER_008", "Unique bridge ID", "String"],
        ["DECK_COND_058", "Deck condition (0-9)", "Integer"],
        ["SUPERSTRUCTURE_COND_059", "Superstructure condition", "Integer"],
        ["SUBSTRUCTURE_COND_060", "Substructure condition", "Integer"],
        ["CULVERT_COND_062", "Culvert condition", "Integer"],
        ["YEAR_BUILT_027", "Year of construction", "Integer"],
        ["ADT_029", "Average daily traffic", "Integer"],
        ["LAT_016 / LONG_017", "Coordinates (DDMMSSCC)", "Packed int"],
    ], Inches(6.8), Inches(1.5), Inches(6.2), Inches(4.0),
       col_widths=[Inches(2.8), Inches(2.2), Inches(1.2)])
    _snum(s, 5)
    _notes(s,
        "The National Bridge Inventory is one of the most comprehensive infrastructure "
        "datasets in the world. Every state, federal agency, and tribal government submits "
        "annual inspection data to the FHWA. Each record captures 123 fields covering "
        "identification, location, structural details, condition ratings, and traffic. The "
        "condition ratings use a zero-to-nine scale where nine is excellent and zero means "
        "the component has failed. Coordinates are encoded in a packed integer format that "
        "we need to convert to decimal degrees. All key analysis columns are 100 percent "
        "populated in the Texas file, making it ideal for teaching.")

    # ── Slide 6: Condition Risk ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Risk Component 1: Structural Condition")
    _rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.1), LIGHT_GREY)
    _tb(s, Inches(1.2), Inches(1.6), Inches(11), Inches(0.9),
        "condition_risk  =  ( 9  -  min_rating )  /  9",
        sz=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, font="Consolas")
    _bullets(s, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.0), [
        "min_rating = minimum of deck, super, sub, culvert",
        "Inverted scale: rating 9 -> risk 0.0 (best)",
        "Rating 0 -> risk 1.0 (worst / failed)",
        "'N' values (not applicable) treated as NaN",
    ], sz=18, sp=Pt(12))
    _table(s, [
        ["Rating", "Meaning", "Risk"],
        ["9", "Excellent", "0.00"],
        ["7", "Good", "0.22"],
        ["5", "Fair", "0.44"],
        ["4", "Poor (FHWA threshold)", "0.56"],
        ["2", "Critical", "0.78"],
        ["0", "Failed", "1.00"],
    ], Inches(7.2), Inches(3.0), Inches(5.5), Inches(3.0),
       col_widths=[Inches(1.0), Inches(2.8), Inches(1.7)])
    _snum(s, 6)
    _notes(s,
        "The first risk component captures current structural condition. The NBI rates four "
        "bridge components -- deck, superstructure, substructure, and culvert -- on a "
        "zero-to-nine scale. We take the minimum across all applicable ratings, because the "
        "weakest component governs structural adequacy. Then we invert the scale so that "
        "higher values mean higher risk. A bridge rated 9 in excellent condition gets "
        "condition risk of zero. One rated 4, the FHWA Poor threshold, gets 0.56. This "
        "linear inversion is simple but effective -- it preserves the ordinality of the NBI "
        "scale while mapping it to a consistent zero-to-one risk range.")

    # ── Slide 7: Age Risk ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Risk Component 2: Age Beyond Design Life")
    _rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.1), LIGHT_GREY)
    _tb(s, Inches(1.2), Inches(1.6), Inches(11), Inches(0.9),
        "age_risk  =  min( age / 50,   1.0 )",
        sz=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, font="Consolas")
    _bullets(s, Inches(0.8), Inches(3.0), Inches(11.5), Inches(3.5), [
        "age = 2025 - year_built",
        "Design life = 50 years (standard assumption)",
        "Linear increase from 0 at age 0 to 1.0 at age 50",
        "Capped at 1.0 -- a 100-year-old bridge is not 'twice as risky' as a 50-year-old",
        "43.8% of Texas bridges exceed 50-year design life",
        "Example: built 1970 -> age 55 -> age_risk = min(55/50, 1.0) = 1.0",
    ], sz=18, sp=Pt(12))
    _snum(s, 7)
    _notes(s,
        "The second component captures age-related deterioration. We normalise bridge age "
        "against a 50-year design life, which is the standard assumption in U.S. bridge "
        "engineering. The key design choice is the cap at 1.0. Without capping, a "
        "125-year-old bridge would get an age risk of 2.5, which would dominate the composite "
        "score and make the other components irrelevant. The cap says: once a bridge has "
        "exceeded its design life, it is at maximum age risk regardless of how far past 50 "
        "years it is. This is a deliberate modelling choice -- additional age beyond the "
        "design life is already reflected in the condition ratings from ongoing inspections.")

    # ── Slide 8: Traffic Risk ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Risk Component 3: Traffic Exposure")
    _rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.1), LIGHT_GREY)
    _tb(s, Inches(1.2), Inches(1.6), Inches(11), Inches(0.9),
        "traffic_risk  =  log1p( ADT )  /  log1p( max_ADT )",
        sz=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, font="Consolas")
    _bullets(s, Inches(0.8), Inches(3.0), Inches(5.5), Inches(3.5), [
        "ADT = average daily traffic (vehicles/day)",
        "Texas range: 0 to 810,110",
        "log1p compresses 5 orders of magnitude",
        "log1p(0) = 0: unused bridges get zero traffic risk",
        "Captures consequence of failure, not structural state",
    ], sz=18, sp=Pt(12))
    _table(s, [
        ["ADT", "Context", "Traffic Risk"],
        ["10", "Rural ranch road", "0.017"],
        ["500", "Farm-to-market road", "0.046"],
        ["5,000", "Urban arterial", "0.063"],
        ["50,000", "Urban freeway", "0.079"],
        ["200,000", "Major interstate", "0.090"],
        ["810,110", "IH-35 (max in TX)", "0.100"],
    ], Inches(7.0), Inches(3.0), Inches(5.8), Inches(3.0),
       col_widths=[Inches(1.3), Inches(2.5), Inches(2.0)])
    _snum(s, 8)
    _notes(s,
        "The third component captures traffic exposure as a proxy for consequence of failure. "
        "A bridge in poor condition on a quiet road is less urgent than one carrying 200,000 "
        "vehicles per day. We use log-one-plus rather than raw ADT because traffic spans "
        "five orders of magnitude. Without the log transform, linear normalisation would "
        "compress 95 percent of bridges near zero and only differentiate the top few percent. "
        "The log transform gives meaningful separation across the full range. Note that "
        "traffic risk is not a structural measure -- it captures how many people are affected "
        "if something goes wrong, which is essential for prioritisation.")

    # ── Slide 9: Composite Risk Score ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Composite Risk Score")
    _rect(s, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.3), LIGHT_GREY)
    _tb(s, Inches(1.0), Inches(1.55), Inches(11.3), Inches(1.1),
        "risk_score  =  0.30 x condition  +  0.30 x age  +  0.40 x traffic",
        sz=26, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER, font="Consolas")
    _bullets(s, Inches(0.8), Inches(3.2), Inches(11.5), Inches(3.5), [
        "All sub-scores normalised to [0, 1] before combining",
        "Traffic weighted highest (0.40): emphasises public safety consequence",
        "Condition and age receive equal weight at 0.30 each",
        "Composite score bounded: 0.0 (no risk) to 1.0 (maximum risk)",
        "Weights are adjustable -- you will experiment with alternatives",
    ], sz=19, sp=Pt(14))
    _snum(s, 9)
    _notes(s,
        "The composite score is a simple weighted sum of the three sub-scores. Traffic "
        "receives the highest weight at 0.40 because it captures the public safety consequence "
        "-- a failing bridge carrying 200,000 vehicles per day is a far greater concern than "
        "one carrying 200. Condition and age each get 0.30. This weighting is a modelling "
        "choice, not a physical law. Different agencies might weight condition more heavily if "
        "their priority is structural health, or weight age more if they are planning "
        "long-term capital budgets. In the hands-on session, you will experiment with "
        "different weight combinations and see how the risk landscape changes.")

    # ── Slide 10: Risk Tiers ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Risk Tier Classification")
    _table(s, [
        ["Tier", "Score Range", "Texas Count", "Percentage", "Action Level"],
        ["Low", "0.00 - 0.24", "1,548", "2.7%", "Routine monitoring"],
        ["Moderate", "0.25 - 0.49", "20,127", "35.3%", "Scheduled inspection"],
        ["High", "0.50 - 0.74", "34,966", "61.4%", "Priority assessment"],
        ["Critical", "0.75 - 1.00", "310", "0.5%", "Immediate review"],
    ], Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.8),
       col_widths=[Inches(1.4), Inches(1.6), Inches(2.0), Inches(1.8), Inches(3.5)])
    _bullets(s, Inches(0.8), Inches(5.2), Inches(11.5), Inches(1.5), [
        "61.9% of Texas bridges score High or Critical",
        "Only 1.2% are rated 'Poor' by condition alone -- the composite reveals latent risk",
    ], sz=19, sp=Pt(12))
    _snum(s, 10)
    _notes(s,
        "We classify the continuous risk score into four actionable tiers. The thresholds "
        "at 0.25, 0.50, and 0.75 divide the zero-to-one range into equal quarters. The "
        "Texas results show that 61.4 percent of bridges fall in the High tier, with the "
        "bulk of the inventory concentrated there. Only 0.5 percent reach Critical, but "
        "61.9 percent are High or Critical -- far more than the 1.2 percent that condition "
        "ratings alone flag as Poor. The discrepancy is the key finding: condition ratings "
        "understate risk because they do not account for age and traffic exposure.")

    # ── Slide 11: Results Summary ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Results Summary")
    _table(s, [
        ["Metric", "Value"],
        ["Total bridges analysed", "56,951"],
        ["Mean risk score", "0.52"],
        ["Median risk score", "0.55"],
        ["Bridges rated Poor (min_cond <= 4)", "680 (1.2%)"],
        ["Bridges exceeding 50-year design life", "24,951 (43.8%)"],
        ["High-traffic Poor bridges (ADT > 50k)", "18"],
        ["Oldest bridge", "1900 (County 217, 125 years)"],
        ["Highest risk score", "0.826 (County 245, built 1962)"],
    ], Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5),
       col_widths=[Inches(5.5), Inches(4.8)])
    _snum(s, 11)
    _notes(s,
        "Here are the headline numbers. The mean risk score of 0.52 places the average Texas "
        "bridge in the High tier. The median of 0.55 is close behind. With traffic weighted "
        "at 40 percent, the model shifts risk toward high-ADT urban corridors. The gap between "
        "the 1.2 percent rated Poor by condition alone and the 61.9 percent scoring High or "
        "Critical by the composite is the central finding. Eighteen bridges in Poor condition "
        "carry over 50,000 vehicles per day, mostly on IH-45 and IH-35E in the Houston and "
        "Dallas-Fort Worth metro areas. These represent the highest-consequence structures.")

    # ── Slide 12: Risk Map ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Geographic Distribution of Risk")
    s.shapes.add_picture(str(FIG1), Inches(0.3), Inches(1.25), Inches(9.5))
    _bullets(s, Inches(10.1), Inches(1.8), Inches(3.0), Inches(4.5), [
        "Clear east-west gradient",
        "Highest risk in rural East Texas",
        "Older road networks (1930s-1960s)",
        "Metro areas show mixed patterns",
        "Western TX: sparser, newer infrastructure",
    ], sz=15, sp=Pt(12))
    _tb(s, Inches(0.3), Inches(6.7), Inches(9.5), Inches(0.5),
        "Figure 1. Composite risk scores for 56,951 Texas structures with metro insets.",
        sz=11, color=MED_GREY)
    _snum(s, 12)
    _notes(s,
        "The risk map reveals a pronounced east-to-west gradient. The highest-risk structures "
        "concentrate in eastern Texas, east of the I-35 corridor, in regions like the Piney "
        "Woods and Gulf Coast. These areas contain the state's oldest road networks, many "
        "built during federal highway expansion in the 1930s through 1960s. Rural "
        "farm-to-market road bridges in counties like Polk and Limestone show particularly "
        "high risk. The metro insets show that even within cities, legacy interstate segments "
        "from the 1950s appear as high-risk clusters embedded in otherwise lower-risk networks.")

    # ── Slide 13: Age vs Condition ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Key Finding: Age Does Not Predict Condition")
    s.shapes.add_picture(str(FIG2), Inches(0.3), Inches(1.25), Inches(8.5))
    _bullets(s, Inches(9.2), Inches(1.8), Inches(3.8), Inches(4.5), [
        "43.8% exceed 50-year design life",
        "Only 1.5% of aged bridges rate as Poor",
        "Dense clustering at ratings 6-7 across all ages",
        "Maintenance sustains condition -- but for how long?",
        "Upper-right quadrant = latent risk zone",
    ], sz=15, sp=Pt(12))
    _tb(s, Inches(0.3), Inches(6.7), Inches(9.0), Inches(0.5),
        "Figure 2. Bridge age vs. min condition rating. Dashed lines: FHWA Poor threshold (4) and 50-year design life.",
        sz=11, color=MED_GREY)
    _snum(s, 13)
    _notes(s,
        "This is arguably the most important figure. It shows that bridge age does not "
        "strongly predict condition in the Texas inventory. Nearly 44 percent of bridges "
        "have exceeded their 50-year design life, but only 1.5 percent of those are rated "
        "Poor. Most cluster at ratings 6 and 7, which is satisfactory, even at ages of 80 "
        "or 100 years. This suggests that sustained maintenance has been effective at "
        "preserving condition. However, this also creates a latent risk zone: these bridges "
        "are acceptable today but vulnerable if maintenance budgets decline. The composite "
        "score captures this risk; condition ratings alone do not.")

    # ── Slide 14: County Risk ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "County-Level Risk Patterns")
    s.shapes.add_picture(str(FIG4), Inches(0.3), Inches(1.25), Inches(8.0))
    _bullets(s, Inches(8.6), Inches(1.8), Inches(4.4), Inches(4.5), [
        "Top 20 counties by mean risk (min. 50 bridges)",
        "All exceed the statewide mean of 0.60",
        "County 235 (Irion) highest at 0.77",
        "Predominantly rural counties",
        "Larger counties (Harris, Tarrant) have more total risk but lower mean",
        "Suggests funding disparity: rural vs urban",
    ], sz=15, sp=Pt(12))
    _tb(s, Inches(0.3), Inches(6.7), Inches(8.0), Inches(0.5),
        "Figure 4. Top 20 counties by mean composite risk score, filtered to min. 50 bridges.",
        sz=11, color=MED_GREY)
    _snum(s, 14)
    _notes(s,
        "The county analysis shows that the highest mean risk scores are in rural counties. "
        "All 20 counties in this chart exceed the statewide mean. This is consistent with "
        "the geographic pattern: rural East Texas counties have older bridge stocks and "
        "receive less maintenance investment per structure than urban areas. Large metro "
        "counties like Harris and Tarrant have more bridges in absolute terms but their "
        "mean risk is lower because new construction dilutes the average. This urban-rural "
        "disparity has direct implications for equitable funding allocation.")

    # ── Slide 15: Limitations ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Model Limitations")
    _bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), [
        "Cross-sectional snapshot: no deterioration trajectory or trend analysis",
        "Equal-interval tier boundaries (0.25) are arbitrary; no empirical calibration",
        "Weights are expert-assigned, not empirically derived from failure data",
        "NBI condition ratings are subjective: inter-inspector variability is documented",
        "Traffic data (ADT) may lag actual volumes; does not capture peak-hour loading",
        "No integration with maintenance expenditure or material-specific decay models",
        "Geographic analysis uses county aggregation; sub-county patterns may differ",
    ], sz=18, sp=Pt(14))
    _snum(s, 15)
    _notes(s,
        "Every model has limitations, and it is important to be transparent about them. "
        "This is a single-year snapshot; we cannot track how risk changes over time without "
        "multi-year data. The tier boundaries and weights are modelling choices, not "
        "empirically calibrated thresholds. NBI ratings are assigned by inspectors in the "
        "field and have documented variability between inspectors. Traffic data may be "
        "several years old and does not capture peak-hour conditions. A production system "
        "would need to integrate maintenance records, material-specific deterioration "
        "models, and longitudinal rating trajectories.")

    # ── Slide 16: Discussion Questions ──
    s = prs.slides.add_slide(bl)
    _bg(s, WHITE)
    _header(s, "Discussion Questions")
    _bullets(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5), [
        "Should condition, age, or traffic receive the highest weight? Why?",
        "How would you validate this model against actual bridge failures?",
        "What additional data would improve the risk assessment?",
        "Should rural and urban bridges be scored with the same formula?",
        "How would you communicate these results to a non-technical audience?",
        "What ethical considerations arise when prioritising some communities over others?",
    ], sz=20, sp=Pt(18))
    _snum(s, 16)
    _notes(s,
        "These questions are designed to stimulate critical thinking about modelling choices. "
        "There is no single correct answer to the weighting question -- it depends on agency "
        "priorities. Validation against actual failures is essential but difficult because "
        "bridge collapses are rare events. Additional data could include maintenance records, "
        "seismic zone classification, or climate exposure indices. The urban-rural question "
        "is particularly relevant given the funding disparity we observed. Finally, risk "
        "prioritisation is inherently an ethical exercise: deprioritising one bridge means "
        "accepting higher risk for the community that depends on it.")

    # ── Slide 17: Hands-On Session ──
    s = prs.slides.add_slide(bl)
    _bg(s, DARK_BLUE)
    _tb(s, Inches(0.8), Inches(0.3), Inches(11.5), Inches(0.9),
        "Hands-On: Interactive Risk Explorer", sz=32, bold=True, color=WHITE)

    _rect(s, Inches(0.5), Inches(1.3), Inches(12.3), Inches(5.5), RGBColor(0x14, 0x2D, 0x4A))

    _tb(s, Inches(0.8), Inches(1.4), Inches(5), Inches(0.5),
        "Open bridge_risk_explorer.html in your browser", sz=18, bold=True, color=ACCENT)

    _tb(s, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5),
        "Task 1: Default Weights", sz=20, bold=True, color=WHITE)
    _bullets(s, Inches(1.0), Inches(2.5), Inches(11), Inches(1.0), [
        "Set weights to Condition=30, Age=30, Traffic=40 (default)",
        "Record: total Critical bridges, mean risk score, geographic pattern",
    ], sz=16, color=RGBColor(0xB0, 0xCC, 0xE6), sp=Pt(6))

    _tb(s, Inches(0.8), Inches(3.5), Inches(11.5), Inches(0.5),
        "Task 2: Condition-Heavy Scenario", sz=20, bold=True, color=WHITE)
    _bullets(s, Inches(1.0), Inches(4.0), Inches(11), Inches(1.0), [
        "Set weights to Condition=80, Age=10, Traffic=10",
        "How does the Critical count change? Which counties move in/out of the top 20?",
    ], sz=16, color=RGBColor(0xB0, 0xCC, 0xE6), sp=Pt(6))

    _tb(s, Inches(0.8), Inches(5.0), Inches(11.5), Inches(0.5),
        "Task 3: No-Traffic Comparison", sz=20, bold=True, color=WHITE)
    _bullets(s, Inches(1.0), Inches(5.5), Inches(11), Inches(1.0), [
        "Set Traffic slider to 0 (Condition=60, Age=40, Traffic=0)",
        "Compare to Task 1: do urban or rural bridges benefit more from removing traffic?",
    ], sz=16, color=RGBColor(0xB0, 0xCC, 0xE6), sp=Pt(6))

    _snum(s, 17)
    _notes(s,
        "This is the hands-on portion. Open the bridge risk explorer HTML file in your "
        "browser. It loads all 56,951 bridges and lets you adjust the risk weights in real "
        "time. Task 1 establishes a baseline using the default weights of 30/30/40. Task 2 "
        "asks you to heavily weight condition -- this will reduce the number of Critical "
        "bridges because most bridges have acceptable condition ratings. Task 3 removes "
        "traffic entirely, which should shift the risk pattern away from urban interstates "
        "and toward older rural structures. Compare the three scenarios and prepare to "
        "discuss what weight combination you would recommend to a state transport agency.")

    prs.save(str(PPTX_PATH))
    print(f"Saved lecture presentation: {PPTX_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
