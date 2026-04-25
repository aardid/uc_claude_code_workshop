"""Generate a 10-slide academic conference presentation from the Texas bridge risk analysis."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUTPUT_DIR = Path("bridge_risk_demo") / "outputs"
PPTX_PATH = OUTPUT_DIR / "presentation.pptx"

FIG1 = OUTPUT_DIR / "fig1_risk_map.png"
FIG2 = OUTPUT_DIR / "fig2_age_condition.png"
FIG3 = OUTPUT_DIR / "fig3_risk_distribution.png"
FIG4 = OUTPUT_DIR / "fig4_county_risk.png"

# -- Colour palette --
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
MID_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_ORANGE = RGBColor(0xD4, 0x6B, 0x08)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GREY = RGBColor(0x33, 0x33, 0x33)
MED_GREY = RGBColor(0x66, 0x66, 0x66)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ROW_BG = RGBColor(0xE8, 0xEE, 0xF4)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def _add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bg(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK_GREY, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _set_para(p, text, font_size=18, bold=False, color=DARK_GREY,
              alignment=PP_ALIGN.LEFT, font_name="Calibri", space_after=Pt(6)):
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after


def _add_bullet_frame(slide, left, top, width, height, bullets, font_size=18,
                      color=DARK_GREY, bold=False, spacing=Pt(10)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox


def _add_slide_number(slide, number):
    _add_textbox(
        slide, Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35),
        str(number), font_size=10, color=MED_GREY, alignment=PP_ALIGN.RIGHT,
    )


def _add_section_header(slide, title):
    bar = _add_shape_bg(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.15), DARK_BLUE)
    _add_textbox(
        slide, Inches(0.6), Inches(0.18), Inches(12), Inches(0.8),
        title, font_size=30, bold=True, color=WHITE,
    )


def _add_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # ==================================================================
    # SLIDE 1 — Title
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, DARK_BLUE)

    _add_textbox(
        slide, Inches(1), Inches(1.5), Inches(11.3), Inches(1.5),
        "Composite Risk Scoring of Bridge Infrastructure",
        font_size=38, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide, Inches(1), Inches(3.0), Inches(11.3), Inches(1.0),
        "A Data-Driven Analysis of 56,951 Texas Structures\nUsing the FHWA National Bridge Inventory",
        font_size=22, bold=False, color=RGBColor(0xB0, 0xCC, 0xE6),
        alignment=PP_ALIGN.LEFT,
    )

    # Author line
    _add_textbox(
        slide, Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
        "David Dempsey",
        font_size=20, bold=True, color=WHITE, alignment=PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.5),
        "University of Canterbury",
        font_size=16, bold=False, color=RGBColor(0xB0, 0xCC, 0xE6),
        alignment=PP_ALIGN.LEFT,
    )

    # Conference footer
    _add_shape_bg(slide, Inches(0), Inches(6.7), SLIDE_WIDTH, Inches(0.8),
                  RGBColor(0x14, 0x2D, 0x4A))
    _add_textbox(
        slide, Inches(1), Inches(6.78), Inches(11.3), Inches(0.5),
        "Data Analysis Workshop  |  2025  |  Civil & Natural Resources Engineering",
        font_size=14, color=RGBColor(0x90, 0xAA, 0xCC), alignment=PP_ALIGN.LEFT,
    )

    _add_notes(slide,
        "Welcome to this presentation on composite risk scoring of bridge infrastructure. "
        "This work applies a transparent, reproducible risk framework to the complete 2025 "
        "Texas National Bridge Inventory dataset, encompassing nearly 57,000 structures. "
        "The motivation is straightforward: bridge agencies need to prioritise inspections "
        "and capital investment, but current condition ratings alone do not capture exposure "
        "to age-related deterioration or traffic loading. We propose a weighted composite "
        "score that combines these three factors into a single metric, and we demonstrate "
        "its application at the state scale using entirely open data."
    )

    # ==================================================================
    # SLIDE 2 — Problem Statement
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Problem Statement")

    bullets = [
        "42% of U.S. bridges are over 50 years old; 7.5% are structurally deficient (ASCE, 2021)",
        "Texas maintains 56,951 NBI structures — the largest state inventory in the U.S.",
        "Annual maintenance budgets cover less than 40% of estimated rehabilitation needs",
        "No standardised, reproducible framework integrates condition, age, and traffic risk",
    ]
    _add_bullet_frame(
        slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5),
        bullets, font_size=22, color=DARK_GREY, spacing=Pt(20),
    )

    _add_textbox(
        slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.8),
        "Source: ASCE 2021 Infrastructure Report Card; FHWA NBI 2025; TxDOT budget reports",
        font_size=12, color=MED_GREY, alignment=PP_ALIGN.LEFT,
    )
    _add_slide_number(slide, 2)

    _add_notes(slide,
        "The American Society of Civil Engineers gives U.S. bridge infrastructure a C grade. "
        "Nationally, 42 percent of all bridges are over 50 years old — exceeding the typical "
        "design life assumed during original construction. Texas carries the largest bridge "
        "inventory of any state, with nearly 57,000 structures. The challenge is not just volume: "
        "annual state budgets for bridge maintenance and rehabilitation consistently fall short "
        "of estimated needs. What is missing is a transparent, data-driven framework that "
        "combines structural condition with age-based deterioration risk and traffic exposure "
        "into a single prioritisation metric. That is the gap this work addresses."
    )

    # ==================================================================
    # SLIDE 3 — Research Question
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Research Question")

    _add_textbox(
        slide, Inches(1), Inches(1.8), Inches(11.3), Inches(1.5),
        "Which Texas bridges pose the greatest structural risk when accounting "
        "for physical condition, age beyond design life, and traffic exposure?",
        font_size=28, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.LEFT,
    )

    sub_questions = [
        "1.  How does a composite risk score distribute across 56,951 structures?",
        "2.  Where are the highest-risk bridges geographically concentrated?",
        "3.  What is the relationship between bridge age and current condition?",
    ]
    _add_bullet_frame(
        slide, Inches(1.2), Inches(3.8), Inches(11), Inches(2.5),
        sub_questions, font_size=20, color=DARK_GREY, spacing=Pt(16),
    )

    _add_slide_number(slide, 3)

    _add_notes(slide,
        "Our central research question asks which bridges should be prioritised for "
        "inspection or intervention when we move beyond simple condition ratings and "
        "incorporate age-based deterioration and traffic loading. The three sub-questions "
        "structure the analysis: first, we examine the overall distribution of composite "
        "risk scores across the inventory; second, we map the geographic concentration "
        "of high-risk structures; third, we investigate whether age reliably predicts "
        "condition, or whether the relationship is more complex. A key design goal is "
        "full reproducibility — the analysis uses only publicly available NBI data and "
        "open-source tools, so any researcher can replicate and extend it."
    )

    # ==================================================================
    # SLIDE 4 — Data
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Data")

    data_bullets = [
        "Source: FHWA National Bridge Inventory, 2025 release",
        "Scope: Texas (state code 48) — 56,951 structures",
        "Format: 123 columns per record; comma-delimited CSV",
        "Key variables: condition ratings, age, traffic, coordinates",
        "Completeness: all analysis columns 100% populated; zero records removed",
    ]
    _add_bullet_frame(
        slide, Inches(0.8), Inches(1.5), Inches(6.0), Inches(3.5),
        data_bullets, font_size=18, color=DARK_GREY, spacing=Pt(12),
    )

    # Summary table
    table_data = [
        ["NBI Column", "Description", "Values"],
        ["DECK_COND_058", "Deck condition rating", "0–9 (9 = excellent)"],
        ["SUPERSTRUCTURE_COND_059", "Superstructure rating", "0–9"],
        ["SUBSTRUCTURE_COND_060", "Substructure rating", "0–9"],
        ["ADT_029", "Average daily traffic", "0–810,110"],
        ["YEAR_BUILT_027", "Year of construction", "1900–2025"],
    ]

    rows, cols = len(table_data), len(table_data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, Inches(7.2), Inches(1.5),
                                        Inches(5.8), Inches(3.2))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(2.6)
    tbl.columns[1].width = Inches(1.9)
    tbl.columns[2].width = Inches(1.3)

    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = table_data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GREY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG

    _add_slide_number(slide, 4)

    _add_notes(slide,
        "The dataset comes from the FHWA National Bridge Inventory, which is publicly "
        "available and updated annually. The 2025 Texas file contains 56,951 records "
        "with 123 columns each. For this analysis, we use 12 key columns: structural "
        "condition ratings for deck, superstructure, substructure, and culvert components; "
        "average daily traffic; year built; geographic coordinates; and county codes. "
        "Importantly, all analysis columns are 100 percent populated — no records were "
        "removed during cleaning. This is a significant advantage of the NBI: it is one of "
        "the most complete infrastructure datasets available for research."
    )

    # ==================================================================
    # SLIDE 5 — Methods (Risk Score)
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Methods — Composite Risk Score")

    # Formula box
    formula_box = _add_shape_bg(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.3),
                                 LIGHT_GREY)
    _add_textbox(
        slide, Inches(1.2), Inches(1.6), Inches(11), Inches(1.1),
        "risk_score  =  0.30 × condition_risk  +  0.30 × age_risk  +  0.40 × traffic_risk",
        font_size=24, bold=True, color=DARK_BLUE, alignment=PP_ALIGN.CENTER,
        font_name="Consolas",
    )

    # Component explanations
    components = [
        "Condition risk  =  (9 − min_rating) / 9        Inverted NBI scale; 0 = best, 1 = worst",
        "Age risk  =  min(age / 50,  1.0)                   Capped at design life of 50 years",
        "Traffic risk  =  log₁ₚ(ADT) / log₁ₚ(max ADT)     Log-normalised to compress 5 orders of magnitude",
    ]
    _add_bullet_frame(
        slide, Inches(1.0), Inches(3.2), Inches(11.5), Inches(2.2),
        components, font_size=16, color=DARK_GREY, spacing=Pt(14),
    )

    # Tier classification
    _add_textbox(
        slide, Inches(0.8), Inches(5.4), Inches(3.0), Inches(0.5),
        "Risk Tier Classification", font_size=18, bold=True, color=DARK_BLUE,
    )

    tier_data = [
        ["Tier", "Score Range"],
        ["Low", "0.00 – 0.24"],
        ["Moderate", "0.25 – 0.49"],
        ["High", "0.50 – 0.74"],
        ["Critical", "0.75 – 1.00"],
    ]
    tbl_shape = slide.shapes.add_table(5, 2, Inches(0.8), Inches(5.85),
                                        Inches(3.5), Inches(1.35))
    tbl = tbl_shape.table
    tbl.columns[0].width = Inches(1.5)
    tbl.columns[1].width = Inches(2.0)
    for r in range(5):
        for c in range(2):
            cell = tbl.cell(r, c)
            cell.text = tier_data[r][c]
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(13)
                paragraph.font.name = "Calibri"
                if r == 0:
                    paragraph.font.bold = True
                    paragraph.font.color.rgb = WHITE
                else:
                    paragraph.font.color.rgb = DARK_GREY
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_HEADER_BG
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ROW_BG

    _add_slide_number(slide, 5)

    _add_notes(slide,
        "The composite risk score combines three sub-scores, each normalised to the zero-to-one "
        "range. Condition risk inverts the NBI rating scale so that lower condition produces "
        "higher risk. Age risk scales linearly with bridge age relative to a 50-year design "
        "life and is capped at one to avoid unbounded scores for the oldest structures. "
        "Traffic risk uses a log-one-plus transformation because average daily traffic spans "
        "five orders of magnitude — from zero to over 800,000. The weighting assigns 30 percent "
        "each to condition and age, with traffic receiving the highest weight at 40 percent "
        "to emphasise public safety consequence. The four tier labels provide an "
        "intuitive classification for engineering decision-making."
    )

    # ==================================================================
    # SLIDE 6 — Results: Geographic Distribution
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Results — Geographic Distribution")

    slide.shapes.add_picture(str(FIG1), Inches(0.4), Inches(1.3), Inches(9.5))

    insights = [
        "East–west gradient: highest risk concentrated in eastern third of the state",
        "Urban interstate corridors show isolated high-risk clusters from 1950s–60s construction",
    ]
    _add_bullet_frame(
        slide, Inches(10.2), Inches(2.0), Inches(2.8), Inches(4.0),
        insights, font_size=15, color=DARK_GREY, spacing=Pt(16),
    )

    _add_textbox(
        slide, Inches(0.4), Inches(6.7), Inches(9.5), Inches(0.5),
        "Figure 1. Composite risk scores for 56,951 Texas NBI structures, 2025. "
        "Insets show Dallas–Fort Worth, Houston, and San Antonio–Austin metro areas.",
        font_size=11, color=MED_GREY, alignment=PP_ALIGN.LEFT,
    )

    _add_slide_number(slide, 6)

    _add_notes(slide,
        "The risk map reveals a pronounced east-to-west gradient. The highest-risk structures "
        "are concentrated in eastern Texas, roughly east of the I-35 corridor, spanning the "
        "Piney Woods and Gulf Coast regions. These areas contain the state's oldest road "
        "networks, many built during the 1930s through 1960s federal highway expansion. Rural "
        "farm-to-market road bridges in counties such as Polk, Houston, and Limestone show "
        "particularly high risk. The western half of the state is predominantly lower risk, "
        "reflecting its sparser and newer infrastructure. Metro insets show that even within "
        "cities, legacy interstate segments from the 1950s appear as high-risk clusters."
    )

    # ==================================================================
    # SLIDE 7 — Results: Age vs Condition
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Results — Age vs. Structural Condition")

    slide.shapes.add_picture(str(FIG2), Inches(0.4), Inches(1.3), Inches(8.5))

    insights = [
        "43.8% of bridges exceed 50-year design life, yet only 1.5% of those rate as Poor",
        "Dense clustering at rating 6–7 across all ages suggests maintenance sustains condition",
    ]
    _add_bullet_frame(
        slide, Inches(9.3), Inches(2.0), Inches(3.7), Inches(4.0),
        insights, font_size=15, color=DARK_GREY, spacing=Pt(16),
    )

    _add_textbox(
        slide, Inches(0.4), Inches(6.7), Inches(9.0), Inches(0.5),
        "Figure 2. Bridge age vs. minimum condition rating, coloured by composite risk score. "
        "Dashed lines mark the FHWA Poor threshold (rating 4) and the 50-year design life.",
        font_size=11, color=MED_GREY, alignment=PP_ALIGN.LEFT,
    )

    _add_slide_number(slide, 7)

    _add_notes(slide,
        "This figure encodes three variables simultaneously: age on the horizontal axis, "
        "minimum condition rating on the vertical axis, and composite risk score as point "
        "colour. The key finding is counterintuitive: age does not strongly predict condition. "
        "Nearly 44 percent of bridges have exceeded their 50-year design life, but only 1.5 "
        "percent of those aged structures are rated Poor. The dense clustering at ratings 6 "
        "and 7, even for bridges over 80 years old, suggests that sustained maintenance and "
        "rehabilitation programmes have been effective at preserving condition. However, this "
        "also means these bridges remain in the latent risk zone — acceptable today, but "
        "potentially vulnerable if maintenance budgets decline."
    )

    # ==================================================================
    # SLIDE 8 — Results: Distribution & Counties
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Results — Risk Distribution & County Analysis")

    slide.shapes.add_picture(str(FIG3), Inches(0.2), Inches(1.3), Inches(6.5))
    slide.shapes.add_picture(str(FIG4), Inches(6.8), Inches(1.3), Inches(6.3))

    _add_textbox(
        slide, Inches(0.2), Inches(6.5), Inches(6.5), Inches(0.7),
        "Figure 3. Distribution of risk scores. 35.8% of bridges "
        "fall in the Critical tier (>0.75). Mean = 0.60, Median = 0.68.",
        font_size=11, color=MED_GREY, alignment=PP_ALIGN.LEFT,
    )

    _add_textbox(
        slide, Inches(6.8), Inches(6.5), Inches(6.3), Inches(0.7),
        "Figure 4. Top 20 counties by mean risk score (min. 50 bridges). "
        "County 235 (Irion) has the highest mean risk at 0.77.",
        font_size=11, color=MED_GREY, alignment=PP_ALIGN.LEFT,
    )

    _add_slide_number(slide, 8)

    _add_notes(slide,
        "The risk distribution is right-skewed, with a pronounced spike near the Critical "
        "threshold of 0.75. Over 20,000 bridges, 35.8 percent of the inventory, score in "
        "the Critical tier. The median of 0.68 exceeds the mean of 0.60, confirming the "
        "left skew. The county analysis shows that all top-20 counties exceed the statewide "
        "mean, with mean risk scores ranging from 0.72 to 0.77. These are predominantly "
        "rural counties with older bridge stocks and lower maintenance investment per "
        "structure. The county chart filters to counties with at least 50 bridges to "
        "avoid small-sample artifacts."
    )

    # ==================================================================
    # SLIDE 9 — Discussion
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Discussion")

    disc_bullets = [
        "Condition data alone understates risk: 1.2% rated Poor, but 66.9% score High or Critical",
        "Composite scoring enables objective prioritisation across a 57,000-structure portfolio",
        "Open NBI data and transparent methods support full reproducibility by other agencies",
        "Geographic targeting of East Texas could address disproportionate share of risk",
        "Limitation: cross-sectional snapshot; longitudinal tracking needed to validate trajectories",
    ]
    _add_bullet_frame(
        slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(4.5),
        disc_bullets, font_size=20, color=DARK_GREY, spacing=Pt(22),
    )

    _add_slide_number(slide, 9)

    _add_notes(slide,
        "Five key discussion points emerge. First, relying solely on FHWA condition ratings "
        "dramatically underestimates risk exposure: only 1.2 percent of bridges are rated "
        "Poor, but nearly 67 percent score High or Critical when age and traffic are included. "
        "Second, the composite framework provides a defensible basis for capital allocation "
        "across a very large portfolio. Third, the entire analysis is built on publicly "
        "available data and open-source code, making it fully replicable. Fourth, the strong "
        "east-west geographic gradient suggests that regionally targeted programmes could "
        "yield outsized returns. The main limitation is that this is a single-year snapshot; "
        "year-over-year trend analysis would be needed to validate deterioration projections."
    )

    # ==================================================================
    # SLIDE 10 — Conclusions
    # ==================================================================
    slide = prs.slides.add_slide(blank_layout)
    _add_background(slide, WHITE)
    _add_section_header(slide, "Conclusions & Future Work")

    _add_textbox(
        slide, Inches(0.8), Inches(1.4), Inches(5.5), Inches(0.5),
        "Conclusions", font_size=22, bold=True, color=DARK_BLUE,
    )

    conclusions = [
        "66.9% of Texas bridges score High or Critical under composite risk assessment, "
        "far exceeding the 1.2% identified by condition ratings alone",
        "Bridge age is the dominant risk driver: 43.8% of structures exceed their 50-year "
        "design life, with risk concentrated in rural East Texas",
        "18 Poor-condition bridges carry over 50,000 vehicles/day on major interstates, "
        "representing the highest consequence-of-failure concentration",
    ]
    _add_bullet_frame(
        slide, Inches(0.8), Inches(1.9), Inches(11.5), Inches(2.3),
        conclusions, font_size=17, color=DARK_GREY, spacing=Pt(14),
    )

    _add_textbox(
        slide, Inches(0.8), Inches(4.2), Inches(5.5), Inches(0.5),
        "Future Work", font_size=22, bold=True, color=DARK_BLUE,
    )
    _add_bullet_frame(
        slide, Inches(0.8), Inches(4.7), Inches(11.5), Inches(0.7),
        ["Longitudinal analysis of NBI rating trajectories linked to maintenance expenditure "
         "records to distinguish durable structural health from deferred risk"],
        font_size=17, color=DARK_GREY, spacing=Pt(8),
    )

    _add_textbox(
        slide, Inches(0.8), Inches(5.5), Inches(5.5), Inches(0.5),
        "References", font_size=18, bold=True, color=DARK_BLUE,
    )
    refs = [
        "ASCE. (2021). 2021 Report Card for America’s Infrastructure. American Society of Civil Engineers.",
        "FHWA. (2025). National Bridge Inventory ASCII Data. Federal Highway Administration. "
        "https://www.fhwa.dot.gov/bridge/nbi/ascii.cfm",
        "Bolukbasi, M., Mohammadi, J., & Arditi, D. (2004). Estimating the future condition of highway "
        "bridge components. Journal of Infrastructure Systems, 10(3), 118–125.",
        "Chase, S. B., Adu-Gyamfi, Y., Aktan, A. E., & Minaie, E. (2016). Synthesis of National and "
        "International Methodologies Used for Bridge Health Indices. FHWA Report No. FHWA-HRT-15-081.",
    ]
    _add_bullet_frame(
        slide, Inches(0.8), Inches(5.9), Inches(11.5), Inches(1.5),
        refs, font_size=11, color=MED_GREY, spacing=Pt(6),
    )

    _add_slide_number(slide, 10)

    _add_notes(slide,
        "To summarise: the composite risk framework reveals that the scale of structural "
        "risk in the Texas bridge inventory is substantially larger than condition ratings "
        "alone suggest. Two-thirds of the inventory warrants attention when age and traffic "
        "are factored in. The geographic concentration in East Texas and on specific "
        "interstate corridors provides actionable targeting information for asset managers. "
        "The most important next step is longitudinal validation: do the bridges we flag as "
        "high risk actually deteriorate faster? Linking NBI rating histories to TxDOT "
        "maintenance expenditure records would answer this question and strengthen the "
        "framework for operational use. Thank you. The code and data are publicly available "
        "for replication."
    )

    # Save
    prs.save(str(PPTX_PATH))
    print(f"Saved presentation: {PPTX_PATH}")


if __name__ == "__main__":
    build_presentation()
